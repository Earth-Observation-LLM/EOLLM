#!/usr/bin/env python3
"""Generate EOLLM dataset presentation as PowerPoint (.pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)       # deep navy
ACCENT  = RGBColor(0x00, 0x96, 0xC7)        # bright blue
ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)        # teal
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xE0, 0xE0, 0xE0)
ORANGE  = RGBColor(0xFF, 0x99, 0x00)
GREEN   = RGBColor(0x00, 0xCC, 0x88)
RED     = RGBColor(0xFF, 0x55, 0x55)
GRAY    = RGBColor(0x99, 0x99, 0x99)
TABLE_HEADER_BG = RGBColor(0x00, 0x50, 0x7A)
TABLE_ROW_BG1   = RGBColor(0x22, 0x22, 0x3A)
TABLE_ROW_BG2   = RGBColor(0x2A, 0x2A, 0x45)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, title, bullets, sub_bullets=None):
    """Add a dark-background slide with title and bullet points."""
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                title, font_size=32, color=ACCENT, bold=True)

    # Accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0

        # Sub-bullets
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sb
                sp.font.size = Pt(16)
                sp.font.color.rgb = GRAY
                sp.font.name = "Calibri"
                sp.space_after = Pt(4)
                sp.level = 1


def add_table(slide, left, top, width, rows_data, col_widths=None):
    """Add a styled table to the slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.4 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT

            # Cell fill
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_BG1 if r % 2 == 1 else TABLE_ROW_BG2

    return table


def add_image_placeholder(slide, left, top, width, height, label="[TODO: Add image]"):
    """Add a rectangle placeholder for an image."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x50)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(height.inches * 20)


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
                "EOLLM", font_size=60, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(1.2),
                "A Multi-City Visual Question Answering Dataset\nfor Urban Understanding",
                font_size=28, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
                "[TODO: Author names]",
                font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.6),
                "[TODO: Affiliation / Conference / Date]",
                font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2: Motivation
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Motivation", [
        "Vision-Language Models (VLMs) show strong performance on natural image QA,\n"
        "but urban/geospatial understanding remains under-evaluated",
        "Existing VQA benchmarks lack multi-modal pairing of satellite + street-level imagery",
        "No standardized benchmark tests both urbanization understanding\n"
        "and cross-view geolocalization in a unified framework",
        "Need: a geographically diverse, metadata-grounded dataset with\n"
        "deterministic (non-subjective) ground truth",
    ])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 3: Dataset at a Glance
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Dataset at a Glance", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    stats = [
        ("13", "Cities across 8 countries"),
        ("1,459", "Unique locations"),
        ("25,696", "Total questions"),
        ("13", "Question types"),
        ("2", "Task families: Urbanization + Geolocalization"),
        ("51K+", "Images (satellite + street view + composites)"),
    ]
    for i, (num, desc) in enumerate(stats):
        row = i // 2
        col = i % 2
        x = Inches(1.2 + col * 5.8)
        y = Inches(1.8 + row * 1.6)

        add_textbox(slide, x, y, Inches(1.8), Inches(0.7),
                    num, font_size=44, color=ORANGE, bold=True)
        add_textbox(slide, x + Inches(2.0), y + Inches(0.1), Inches(3.5), Inches(0.6),
                    desc, font_size=20, color=LIGHT)

    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                "[TODO: Update numbers with final dataset statistics]",
                font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 4: Pipeline Overview
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Pipeline Overview", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    steps = [
        ("1", "Sample\nLocations", "OSM road\nnetwork"),
        ("2", "Fetch\nSatellite", "GEE / ESRI\n/ IGN"),
        ("3", "Fetch\nStreet View", "Google SV\nAPI"),
        ("4", "Enrich\nMetadata", "OSM +\nNominatim"),
        ("5", "Generate\nComposites", "Arrows +\nGrids"),
        ("6", "Generate\nQuestions", "Templates\n+ MCQ"),
        ("7", "Validate", "Quality\nfilters"),
    ]

    for i, (num, label, detail) in enumerate(steps):
        x = Inches(0.5 + i * 1.75)
        y = Inches(2.2)

        # Step box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, y, Inches(1.5), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x45)
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1.5)

        # Step number
        add_textbox(slide, x + Inches(0.05), y + Inches(0.1), Inches(1.4), Inches(0.4),
                    num, font_size=24, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
        # Step label
        add_textbox(slide, x + Inches(0.05), y + Inches(0.5), Inches(1.4), Inches(0.6),
                    label, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Detail
        add_textbox(slide, x + Inches(0.05), y + Inches(1.1), Inches(1.4), Inches(0.6),
                    detail, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

        # Arrow between steps
        if i < len(steps) - 1:
            add_textbox(slide, x + Inches(1.5), y + Inches(0.65), Inches(0.3), Inches(0.4),
                        "\u2192", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.8),
                "Linear chain: each step enriches sample dicts passed to the next.\n"
                "Final output: dataset.jsonl with all questions, images, and validation flags.",
                font_size=16, color=LIGHT)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 5: Geographic Coverage
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Geographic Coverage", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    add_image_placeholder(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5),
                          "[TODO: World map with 13 city markers]")

    cities_text = (
        "Turkey: Samsun, Antalya\n"
        "Switzerland: Zurich\n"
        "Russia: St. Petersburg, Moscow\n"
        "Portugal: Lisbon\n"
        "Cyprus: Nicosia\n"
        "Iceland: Reykjavik\n"
        "Estonia: Tallinn\n"
        "Argentina: Buenos Aires\n"
        "Belgium: Brussels\n"
        "Austria: Vienna\n"
        "Hungary: Budapest"
    )
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(1.5), Inches(5.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line_txt in cities_text.split("\n"):
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = line_txt
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT
        p.font.name = "Calibri"
        p.space_after = Pt(4)

    add_textbox(slide, Inches(7.2), Inches(5.8), Inches(5.5), Inches(0.8),
                "5 urban character types per city:\nCommercial, Residential, Mixed-use, Industrial, Natural/Recreation",
                font_size=14, color=GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 6: Data Sources
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Data Sources", [
        "Satellite imagery (512x512 px) -- multi-source with automatic fallback:",
        "Street View imagery (640x640 px) -- Google Street View Static API",
        "OpenStreetMap metadata -- Overpass API (buildings, roads, amenities, land use, transit, water)",
        "Reverse geocoding -- Nominatim API (suburb, district, postcode)",
    ], sub_bullets={
        0: [
            "NAIP via Google Earth Engine (US) -- 0.6 m/px",
            "IGN Geoplateforme WMS (France) -- 0.2 m/px",
            "ESRI World Imagery REST (Europe, Turkey) -- 0.3-0.5 m/px",
            "Sentinel-2 via GEE (global fallback) -- 10 m/px",
        ],
        1: [
            "4 road-aligned headings: forward, backward, left, right",
            "Outdoor-only, tunnel detection heuristic, snap distance < 80 m",
        ],
    })

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 7: Location Sampling
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Location Sampling Strategy", [
        "11 hand-curated seed locations per city covering 5 urban character types",
        "Random sampling within city bounding box, then road-snapping via OSM",
        "Road bearing extracted for consistent street-view heading alignment",
        "Locations rejected if no OSM road within ~200 m radius",
        "~120 samples per city (configurable)",
    ], sub_bullets={
        0: [
            "Commercial (2), Residential (3), Mixed-use (2), Industrial (2), Natural/Recreation (2)",
        ],
        1: [
            "Snap to nearest road node; extract highway type and bearing angle (0-360 deg)",
        ],
    })

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 8: Image Acquisition
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Image Acquisition", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    # Satellite side
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                "Satellite Tile", font_size=22, color=ORANGE, bold=True)
    add_image_placeholder(slide, Inches(0.8), Inches(2.1), Inches(2.5), Inches(2.5),
                          "[TODO: Sample\nsatellite tile]")
    sat_bullets = [
        "512x512 px, configurable buffer radius",
        "Auto-detected source per region",
        "Fallback chain: NAIP > IGN > ESRI > S2",
        "Quality filter: reject < 5 KB tiles",
        "Cloud filter: S2 max 20% coverage",
    ]
    txBox = slide.shapes.add_textbox(Inches(3.6), Inches(2.1), Inches(3), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for b in sat_bullets:
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = "\u2022 " + b
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT
        p.font.name = "Calibri"
        p.space_after = Pt(4)

    # Street View side
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                "Street View (4 angles)", font_size=22, color=ORANGE, bold=True)
    add_image_placeholder(slide, Inches(7), Inches(2.1), Inches(2.5), Inches(2.5),
                          "[TODO: Sample\n4-angle SV grid]")
    sv_bullets = [
        "640x640 px, FOV 90 deg, pitch -5 deg",
        "4 headings: fwd, bwd, left, right",
        "Road-bearing aligned for consistency",
        "Tunnel detection via RGB analysis",
        "Snap distance validated < 80 m",
    ]
    txBox = slide.shapes.add_textbox(Inches(9.8), Inches(2.1), Inches(3), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for b in sv_bullets:
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = "\u2022 " + b
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT
        p.font.name = "Calibri"
        p.space_after = Pt(4)

    add_image_placeholder(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2),
                          "[TODO: Example showing satellite + 4 street view angles for one location]")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 9: OSM Metadata Extraction
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "OSM Metadata Extraction", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    metadata_table = [
        ["Feature", "Source", "Radius", "Derived Attribute"],
        ["Buildings", "OSM building tags", "200 m", "Count, median levels, type distribution"],
        ["Land Use", "OSM landuse + heuristic", "200 m", "8 categories (residential, commercial, ...)"],
        ["Roads", "OSM highway tags", "At point", "Type, surface material, junction type"],
        ["Amenities", "OSM amenity tags", "200 m", "Count, type list, richness level"],
        ["Green Space", "OSM leisure=park", "200 m", "Boolean presence"],
        ["Transit", "OSM bus/train stops", "300 m", "Stop count, density level"],
        ["Water", "OSM waterway/natural", "Nearest", "Distance in meters, proximity bin"],
    ]
    add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), metadata_table,
              col_widths=[Inches(1.5), Inches(2.8), Inches(1.2), Inches(6.8)])

    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1),
                "Land use classification uses a priority cascade:\n"
                "Explicit OSM tag > Building type distribution > Amenity density ratio > Default (residential)",
                font_size=15, color=GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 10: Urbanization Tasks Overview
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Urbanization Tasks", [
        "10 question types probing urban structure, land use, and infrastructure",
        "Ground truth derived deterministically from OSM metadata (no human annotation)",
        "Template-based generation: 10 natural-language paraphrases per question type",
        "Multiple-choice format (A/B/C/D) with semantically valid distractors",
        "Each sample generates all feasible questions; best one selected via diversity heuristic",
    ])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 11: Urbanization Question Types Table
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8),
                "Urbanization: Question Types", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(0.95), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    q_table = [
        ["Question Type", "Answer Categories", "Difficulty", "Modality"],
        ["Land Use", "8 types (residential, commercial, industrial, ...)", "Easy", "Sat + SV"],
        ["Building Height", "4 ranges (1-3, 4-7, 8-20, 20+ floors)", "Easy", "Sat + SV"],
        ["Urban Density", "4 levels (0-15, 16-50, 51-150, 150+ bldgs)", "Medium", "Sat"],
        ["Road Type", "6 types (motorway to residential)", "Easy", "SV"],
        ["Road Surface", "4 types (asphalt, cobblestone, unpaved, concrete)", "Easy", "SV"],
        ["Junction Type", "4 types (roundabout, signalized, ...)", "Medium", "SV"],
        ["Green Space", "Yes / No + context", "Easy", "Sat + SV"],
        ["Amenity Richness", "4 levels (high / moderate / low / minimal)", "Medium", "SV"],
        ["Transit Density", "4 levels (none / low / moderate / high)", "Medium", "Sat"],
        ["Water Proximity", "3 bins (0-50m, 50-150m, 150m+)", "Medium", "Sat"],
    ]
    add_table(slide, Inches(0.3), Inches(1.2), Inches(12.7), q_table,
              col_widths=[Inches(2.2), Inches(5.5), Inches(1.8), Inches(2.0)])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 12: Answer & Distractor Generation
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Answer & Distractor Generation", [
        "Answers are deterministic: derived programmatically from OSM tags and geometry",
        "Distractors: all other valid options from the answer space",
        "Option order shuffled per question with reproducible RNG (seed=42)",
        "Answer key (A/B/C/D) distribution validated for balance across dataset",
        "Questions skipped if metadata is missing or ambiguous (e.g., no building:levels tag)",
    ], sub_bullets={
        0: [
            "Example: building_height answer = bin(median OSM building:levels within 200m)",
            "Example: land_use = priority cascade over OSM tags, building types, amenity density",
        ],
        1: [
            "e.g., for building_height with answer '4-7 floors', distractors = {'1-3', '8-20', '20+'}",
        ],
    })

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 13: Geolocalization Tasks Overview
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Geolocalization Tasks", [
        "3 cross-modal alignment tasks testing satellite <-> street-view correspondence",
        "Camera Direction: which satellite arrow matches the street-view viewing angle?",
        "Mismatch Binary: does this street-view set match the marked satellite location? (yes/no)",
        "Mismatch MCQ: which of 4 street-view sets corresponds to the marked satellite? (4-way)",
        "Two difficulty strategies: same-city negatives (harder) vs cross-city negatives (easier)",
    ])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 14: Camera Direction Task
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Geolocalization: Camera Direction", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
                "Task: Given a street-view image, identify which satellite image (with directional arrow)\n"
                "shows the correct viewing direction.",
                font_size=18, color=LIGHT)

    # Left: street view placeholder
    add_image_placeholder(slide, Inches(0.8), Inches(2.8), Inches(3.5), Inches(3.5),
                          "[TODO: Street view image\n(e.g., forward direction)]")

    # Right: 4 satellite options
    dirs = ["A: Forward", "B: Right", "C: Backward", "D: Left"]
    for i, d in enumerate(dirs):
        row = i // 2
        col = i % 2
        x = Inches(5.2 + col * 3.8)
        y = Inches(2.8 + row * 1.8)
        add_image_placeholder(slide, x, y, Inches(3.2), Inches(1.5),
                              f"[Satellite + arrow: {d}]")

    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
                "Red directional arrows overlaid at satellite tile center, aligned to road bearing + 90 deg offsets",
                font_size=14, color=GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 15: Mismatch Tasks
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Geolocalization: Mismatch Tasks", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    # Binary section
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                "Mismatch Binary (Yes/No)", font_size=22, color=ORANGE, bold=True)
    binary_bullets = [
        "\u2022 Satellite (red dot) + street-view composite",
        "\u2022 Positive: own location pair (answer: Yes)",
        "\u2022 Negative: different location (answer: No)",
        "\u2022 Same-city negatives = harder",
        "\u2022 Cross-city negatives = easier",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for b in binary_bullets:
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT
        p.font.name = "Calibri"
        p.space_after = Pt(4)

    # MCQ section
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                "Mismatch MCQ (4-Way)", font_size=22, color=ORANGE, bold=True)
    mcq_bullets = [
        "\u2022 Satellite (red dot) + 2x2 grid of SV composites",
        "\u2022 1 correct location + 3 distractors",
        "\u2022 Each cell = 4-angle SV composite",
        "\u2022 Answer: grid position (A/B/C/D)",
        "\u2022 Hardest geolocalization task",
    ]
    txBox = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(5.5), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for b in mcq_bullets:
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT
        p.font.name = "Calibri"
        p.space_after = Pt(4)

    # Example placeholders
    add_image_placeholder(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.5),
                          "[TODO: Mismatch Binary example\n(satellite + SV composite pair)]")
    add_image_placeholder(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(2.5),
                          "[TODO: Mismatch MCQ example\n(satellite + 2x2 SV grid)]")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 16: Question Selection & Diversity
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Question Selection & Diversity", [
        "Each location generates 10-20 feasible questions; all stored in dataset",
        "Best question per location selected via scoring heuristic:",
        "Result: balanced topic distribution across dataset, favoring rare and distinctive questions",
    ], sub_bullets={
        1: [
            "Topic rarity: penalize overused topics (score -= usage x 3)",
            "Visual distinctiveness: bonus for extreme/unusual values (e.g., 20+ floor buildings)",
            "Cross-modal preference: geolocalization tasks get baseline bonus",
            "Rare land-use types rewarded (industrial, institutional > residential)",
        ],
    })

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 17: Quality Assurance
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Quality Assurance & Validation", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    val_table = [
        ["Check", "Level", "Criterion"],
        ["Question presence", "CRITICAL", "Must have question field"],
        ["Answer in options", "CRITICAL", "Answer key exists in A/B/C/D"],
        ["Unique options", "CRITICAL", "All 4 options are distinct"],
        ["Satellite image", "CRITICAL", "File exists AND > 5 KB"],
        ["Street view images", "WARNING", "At least 1 of 4 angles present"],
        ["Answer distribution", "WARNING", "No single answer > 40%"],
        ["Topic diversity", "WARNING", ">=3 distinct topics in dataset"],
        ["City coverage", "WARNING", ">=2 cities represented"],
    ]
    add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), val_table,
              col_widths=[Inches(3.0), Inches(2.0), Inches(7.3)])

    add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1),
                "Additional filters: tunnel detection (RGB channel spread heuristic),\n"
                "Street View snap distance < 80 m, Sentinel-2 cloud coverage < 20%",
                font_size=15, color=GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 18: Dataset Splits
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Dataset Splits", [
        "Three-way location-level split: Train / Validation / Benchmark",
        "Seen cities (training): 32 cities used for training and in-distribution evaluation",
        "Unseen cities (zero-shot): 8 held-out cities for generalization testing",
        "Benchmark: 10% per city, includes both seen and unseen cities",
        "Leak auditing: no distractor locations from unseen cities appear in training",
    ], sub_bullets={
        2: [
            "Istanbul, Moscow, Chicago, Seoul, Rio, Cape Town, Sydney, Singapore",
        ],
        3: [
            "Stratified by question type with per-type downsampling targets",
        ],
    })

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 19: Statistics Placeholder
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Dataset Statistics", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    add_image_placeholder(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5),
                          "[TODO: Question type distribution chart]")
    add_image_placeholder(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5),
                          "[TODO: Difficulty distribution chart]")
    add_image_placeholder(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.5),
                          "[TODO: City sample distribution chart]")
    add_image_placeholder(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(2.5),
                          "[TODO: Answer key balance chart]")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 20: Statistics Placeholder 2
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Dataset Statistics (cont.)", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    add_image_placeholder(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5),
                          "[TODO: Land use distribution]")
    add_image_placeholder(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5),
                          "[TODO: Metadata coverage heatmap]")
    add_image_placeholder(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.5),
                          "[TODO: Additional statistics]")
    add_image_placeholder(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(2.5),
                          "[TODO: Additional statistics]")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 21: Example Samples
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                "Example Samples", font_size=32, color=ACCENT, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.15), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    add_image_placeholder(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5),
                          "[TODO: Urbanization task example\n\nSatellite + Street View + Question + Options]")
    add_image_placeholder(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.5),
                          "[TODO: Geolocalization task example\n\nSatellite (arrow/dot) + SV composite + Question]")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 22: Summary
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, "Summary & Key Contributions", [
        "Multi-modal urban VQA dataset: satellite + street-view imagery across 13 global cities",
        "Two complementary task families: urbanization understanding + cross-view geolocalization",
        "Fully automated pipeline: deterministic ground truth from OSM (no manual annotation)",
        "13 question types with template-based linguistic diversity (10 paraphrases each)",
        "Rigorous quality assurance: multi-level validation, tunnel detection, leak auditing",
        "Benchmark-ready splits with seen/unseen city partitioning for zero-shot evaluation",
    ])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 23: Thank You
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                "Thank You", font_size=48, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                "[TODO: Contact information / GitHub link]",
                font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
                "Questions?",
                font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ── Save ────────────────────────────────────────────────────────────────
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "EOLLM_Dataset_Presentation.pptx")
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
